#!/usr/bin/env python3
"""
주간 주목 기사 슬라이드 생성 스크립트 (python-pptx).

사용법:
  python slide.py <output_path> <input_json_path>

입력 JSON 형식:
  {
    "제목": "기사 제목",
    "매체": "출처 매체명",
    "날짜": "YYYY-MM-DD",
    "주차": "3월 3주차",
    "기사요약": "5~8문장 요약 텍스트",
    "기관소개": [{"기관명": "...", "설명": "..."}],
    "기술소개": [{"기술명": "...", "설명": "..."}],
    "적용아이디어": [{"제목": "...", "설명": "..."}]
  }

출력: .pptx 파일 (3장)
  1. 기사 요약 — 무슨 일이 있었나
  2. 기관 & 기술 — 누가, 어떻게
  3. 적용 아이디어 — 우리한테 뭐냐
"""

import json
import io
import os
import sys

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 디자인 상수 ──
BG_DARK = RGBColor(0x1E, 0x29, 0x3B)      # 짙은 남색
BG_LIGHT = RGBColor(0xF8, 0xFA, 0xFC)     # 밝은 회색
ACCENT = RGBColor(0x25, 0x63, 0xEB)        # 파란색
TEXT_DARK = RGBColor(0x1E, 0x29, 0x3B)
TEXT_MUTED = RGBColor(0x64, 0x74, 0x8B)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
GREEN = RGBColor(0x10, 0xB9, 0x81)

FONT_TITLE = "맑은 고딕"
FONT_BODY = "맑은 고딕"

SLIDE_WIDTH = Inches(13.333)   # 16:9
SLIDE_HEIGHT = Inches(7.5)

TOTAL_SLIDES = 3


def set_slide_bg(slide, color):
    """슬라이드 배경색 설정."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=14,
                font_color=TEXT_DARK, bold=False, alignment=PP_ALIGN.LEFT,
                font_name=FONT_BODY):
    """텍스트 박스 추가."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=13,
                    font_color=TEXT_DARK):
    """불릿 리스트 추가."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"•  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = FONT_BODY
        p.space_after = Pt(8)

    return txBox


def add_accent_bar(slide, left, top, width, height, color=ACCENT):
    """색상 바 추가."""
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, title, body,
             title_color=ACCENT, border_color=ACCENT):
    """카드 스타일 텍스트 블록 (좌측 색상 바 + 제목 + 본문)."""
    add_accent_bar(slide, left, top, Inches(0.06), height, border_color)
    add_textbox(slide, left + Inches(0.25), top + Inches(0.1), width - Inches(0.3), Inches(0.35),
                title, font_size=14, font_color=title_color, bold=True)
    add_textbox(slide, left + Inches(0.25), top + Inches(0.45), width - Inches(0.3), height - Inches(0.55),
                body, font_size=12, font_color=TEXT_DARK)


def add_footer(slide):
    """하단 프로젝트 정보 바."""
    add_accent_bar(slide, Inches(0), Inches(6.9), SLIDE_WIDTH, Inches(0.6), BG_DARK)
    add_textbox(slide, Inches(0.8), Inches(6.95), Inches(11.5), Inches(0.4),
                "NFA 응급환자 최적 이송 플랫폼  |  AI 이송병원 추천 · 실시간 모니터링 · 응급실 상황판",
                font_size=11, font_color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)


def add_slide_number(slide, num, total):
    """우측 하단 페이지 번호."""
    add_textbox(slide, Inches(12.2), Inches(6.95), Inches(1), Inches(0.4),
                f"{num} / {total}", font_size=10, font_color=TEXT_MUTED,
                alignment=PP_ALIGN.RIGHT)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 슬라이드 1: 기사 요약 — 무슨 일이 있었나
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def build_slide1(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_LIGHT)

    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ACCENT)

    # 주차 태그
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(3), Inches(0.4),
                f"📰 {data['주차']} 주간 브리핑", font_size=14,
                font_color=ACCENT, bold=True)

    # 제목
    add_textbox(slide, Inches(0.8), Inches(0.9), Inches(11.5), Inches(1.0),
                data["제목"], font_size=28, font_color=TEXT_DARK, bold=True,
                font_name=FONT_TITLE)

    # 매체 · 날짜
    add_textbox(slide, Inches(0.8), Inches(1.9), Inches(6), Inches(0.4),
                f"{data['매체']}  ·  {data['날짜']}", font_size=13,
                font_color=TEXT_MUTED)

    add_accent_bar(slide, Inches(0.8), Inches(2.4), Inches(1.5), Inches(0.04), ACCENT)

    add_textbox(slide, Inches(0.8), Inches(2.7), Inches(5.5), Inches(0.4),
                "⚡ 핵심 내용", font_size=16, font_color=ORANGE, bold=True)

    # 요약 텍스트를 문장 단위로 불릿 처리
    sentences = [s.strip() for s in data["기사요약"].split(". ") if s.strip()]
    sentences = [s if s.endswith(".") else s + "." for s in sentences]

    add_bullet_list(slide, Inches(0.8), Inches(3.2), Inches(11.5), Inches(3.5),
                    sentences, font_size=14, font_color=TEXT_DARK)

    add_slide_number(slide, 1, TOTAL_SLIDES)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 슬라이드 2: 기관 & 기술 — 누가, 어떻게
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def build_slide2(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_LIGHT)

    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ACCENT)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.5),
                "🏢 주요 기관 & 기술", font_size=24,
                font_color=TEXT_DARK, bold=True, font_name=FONT_TITLE)

    add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11.5), Inches(0.4),
                f"참조: {data['제목']}", font_size=12, font_color=TEXT_MUTED)

    add_accent_bar(slide, Inches(0.8), Inches(1.5), Inches(1.5), Inches(0.04), ACCENT)

    # 기관 소개 (좌측)
    orgs = data.get("기관소개", [])
    card_top = Inches(1.8)
    for i, org in enumerate(orgs[:3]):
        colors = [ACCENT, GREEN, ORANGE]
        add_card(slide, Inches(0.8), card_top, Inches(5.5), Inches(1.5),
                 org.get("기관명", ""), org.get("설명", ""),
                 title_color=colors[i % 3], border_color=colors[i % 3])
        card_top += Inches(1.65)

    # 기술 소개 (우측)
    techs = data.get("기술소개", [])
    card_top = Inches(1.8)
    for i, tech in enumerate(techs[:2]):
        colors = [ORANGE, ACCENT]
        add_card(slide, Inches(6.8), card_top, Inches(5.7), Inches(2.3),
                 tech.get("기술명", ""), tech.get("설명", ""),
                 title_color=colors[i % 2], border_color=colors[i % 2])
        card_top += Inches(2.45)

    add_slide_number(slide, 2, TOTAL_SLIDES)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 슬라이드 3: 적용 아이디어 — 우리한테 뭐냐
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def build_slide3(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_LIGHT)

    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ORANGE)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.5),
                "💡 NFA 플랫폼 적용 아이디어", font_size=24,
                font_color=TEXT_DARK, bold=True, font_name=FONT_TITLE)

    add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11.5), Inches(0.4),
                f"참조: {data['제목']}", font_size=12, font_color=TEXT_MUTED)

    add_accent_bar(slide, Inches(0.8), Inches(1.5), Inches(1.5), Inches(0.04), ORANGE)

    # 적용 아이디어 카드
    ideas = data.get("적용아이디어", [])
    card_top = Inches(1.8)
    colors = [ACCENT, GREEN, ORANGE]

    for i, idea in enumerate(ideas[:3]):
        add_card(slide, Inches(0.8), card_top, Inches(11.5), Inches(1.5),
                 idea.get("제목", ""), idea.get("설명", ""),
                 title_color=colors[i % 3], border_color=colors[i % 3])
        card_top += Inches(1.65)

    add_slide_number(slide, 3, TOTAL_SLIDES)


def main():
    if len(sys.argv) < 3:
        print("Usage: python slide.py <output_path> <input_json_path>")
        sys.exit(1)

    output_path = sys.argv[1]
    input_path = sys.argv[2]

    with open(input_path, "r", encoding="utf-8-sig") as f:
        data = json.load(f)

    output_dir = os.path.dirname(output_path)
    if output_dir:
        os.makedirs(output_dir, exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    build_slide1(prs, data)
    build_slide2(prs, data)
    build_slide3(prs, data)

    prs.save(output_path)
    print(f"DONE:{output_path}")


if __name__ == "__main__":
    main()
