#!/usr/bin/env python3
"""주목 기사 3장 슬라이드 생성기."""

from __future__ import annotations

import io
import json
import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

BG_LIGHT = RGBColor(0xF8, 0xFA, 0xFC)
BG_DARK = RGBColor(0x1E, 0x29, 0x3B)
TEXT_DARK = RGBColor(0x1E, 0x29, 0x3B)
TEXT_MUTED = RGBColor(0x64, 0x74, 0x8B)
ACCENT = RGBColor(0x25, 0x63, 0xEB)
GREEN = RGBColor(0x10, 0xB9, 0x81)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)

FONT_TITLE = "Malgun Gothic"
FONT_BODY = "Malgun Gothic"

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    size=14,
    color=TEXT_DARK,
    bold=False,
    align=PP_ALIGN.LEFT,
    font=FONT_BODY,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    return box


def add_bar(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_bullets(slide, left, top, width, height, items, *, size=14, color=TEXT_DARK):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    for index, item in enumerate(items):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = FONT_BODY
        p.space_after = Pt(8)
    return box


def add_card(slide, left, top, width, height, title, body, color):
    add_bar(slide, left, top, Inches(0.07), height, color)
    add_textbox(
        slide,
        left + Inches(0.25),
        top + Inches(0.08),
        width - Inches(0.3),
        Inches(0.35),
        title,
        size=15,
        color=color,
        bold=True,
    )
    add_textbox(
        slide,
        left + Inches(0.25),
        top + Inches(0.45),
        width - Inches(0.3),
        height - Inches(0.55),
        body,
        size=12,
    )


def add_slide_number(slide, current: int, total: int) -> None:
    add_textbox(
        slide,
        Inches(12.15),
        Inches(7.0),
        Inches(0.8),
        Inches(0.25),
        f"{current} / {total}",
        size=10,
        color=TEXT_MUTED,
        align=PP_ALIGN.RIGHT,
    )


def split_sentences(text: str) -> list[str]:
    raw = [part.strip() for part in text.replace("\n", " ").split(".") if part.strip()]
    return [item + "." for item in raw]


def build_slide1(prs: Presentation, data: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_LIGHT)
    add_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ACCENT)
    add_textbox(
        slide,
        Inches(0.8),
        Inches(0.38),
        Inches(3.2),
        Inches(0.3),
        f"핫 {data['주차']} 주간 브리핑",
        size=14,
        color=ACCENT,
        bold=True,
    )
    add_textbox(slide, Inches(0.8), Inches(0.9), Inches(11.5), Inches(0.9), data["제목"], size=26, bold=True, font=FONT_TITLE)
    add_textbox(slide, Inches(0.8), Inches(1.85), Inches(8), Inches(0.3), f"{data['매체']} · {data['날짜']}", size=12, color=TEXT_MUTED)
    add_bar(slide, Inches(0.8), Inches(2.35), Inches(1.5), Inches(0.04), ACCENT)
    add_textbox(slide, Inches(0.8), Inches(2.7), Inches(4), Inches(0.3), "이번 주 주목 내용", size=16, color=ORANGE, bold=True)
    add_bullets(slide, Inches(0.8), Inches(3.15), Inches(11.3), Inches(3.4), split_sentences(data["기사요약"]), size=14)
    add_slide_number(slide, 1, 3)


def build_slide2(prs: Presentation, data: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_LIGHT)
    add_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ACCENT)
    add_textbox(slide, Inches(0.8), Inches(0.38), Inches(11), Inches(0.4), "주요 기관과 기술", size=24, bold=True, font=FONT_TITLE)
    add_textbox(slide, Inches(0.8), Inches(0.95), Inches(11), Inches(0.25), f"참조: {data['제목']}", size=11, color=TEXT_MUTED)
    add_bar(slide, Inches(0.8), Inches(1.45), Inches(1.5), Inches(0.04), ACCENT)

    for index, org in enumerate(data.get("기관소개", [])[:3]):
        add_card(
            slide,
            Inches(0.8),
            Inches(1.8 + 1.65 * index),
            Inches(5.5),
            Inches(1.5),
            org.get("기관명", ""),
            org.get("설명", ""),
            [ACCENT, GREEN, ORANGE][index % 3],
        )

    for index, tech in enumerate(data.get("기술소개", [])[:2]):
        add_card(
            slide,
            Inches(6.8),
            Inches(1.8 + 2.45 * index),
            Inches(5.7),
            Inches(2.3),
            tech.get("기술명", ""),
            tech.get("설명", ""),
            [ORANGE, ACCENT][index % 2],
        )

    add_slide_number(slide, 2, 3)


def build_slide3(prs: Presentation, data: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_LIGHT)
    add_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ORANGE)
    add_textbox(slide, Inches(0.8), Inches(0.38), Inches(11), Inches(0.4), "시사점", size=24, bold=True, font=FONT_TITLE)
    add_textbox(slide, Inches(0.8), Inches(0.95), Inches(11), Inches(0.25), f"참조: {data['제목']}", size=11, color=TEXT_MUTED)
    add_bar(slide, Inches(0.8), Inches(1.45), Inches(1.5), Inches(0.04), ORANGE)

    for index, item in enumerate(data.get("시사점", [])[:3]):
        add_card(
            slide,
            Inches(0.8),
            Inches(1.8 + 1.65 * index),
            Inches(11.5),
            Inches(1.5),
            item.get("제목", ""),
            item.get("설명", ""),
            [ACCENT, GREEN, ORANGE][index % 3],
        )

    add_bar(slide, Inches(0), Inches(6.9), SLIDE_WIDTH, Inches(0.6), BG_DARK)
    add_textbox(
        slide,
        Inches(0.8),
        Inches(6.98),
        Inches(11.7),
        Inches(0.2),
        "NFA 응급환자 최적 이송 플랫폼 | AI 이송병원 추천 · 실시간 모니터링 · 응급실 상황판",
        size=10,
        color=TEXT_MUTED,
        align=PP_ALIGN.CENTER,
    )
    add_slide_number(slide, 3, 3)


def main() -> None:
    if len(sys.argv) < 3:
        print("Usage: python slide.py <output_path> <input_json_path>")
        sys.exit(1)

    output_path = sys.argv[1]
    input_path = sys.argv[2]

    with open(input_path, "r", encoding="utf-8-sig") as handle:
        data = json.load(handle)

    output_dir = os.path.dirname(output_path)
    if output_dir:
        os.makedirs(output_dir, exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    build_slide1(prs, data)
    build_slide2(prs, data)
    build_slide3(prs, data)

    prs.save(output_path)
    print(f"DONE:{output_path}")


if __name__ == "__main__":
    main()
